import streamlit as st
from pptx import Presentation
import easyocr
import cv2
import numpy as np
import io
from PIL import Image

st.set_page_config(page_title="NBLM 圖片去字工具專業版", layout="centered")
st.title("🖼️ NBLM 簡報圖片去中文字 (強化標點符號版)")

# 側邊欄設定
st.sidebar.header("工具設定")
lang_option = st.sidebar.radio(
    "1. 選擇語言類型：",
    ('繁體中文', '簡體中文'),
    index=0
)

# 新增清除強度設定 (解決標點符號殘留)
clean_strength = st.sidebar.slider(
    "2. 清除邊緣強度 (數值越大，越能清除標點)：",
    min_value=1, max_value=20, value=10, help="這會讓去色範圍向外擴張，確保包裹住細小的標點符號。"
)

@st.cache_resource
def load_ocr(lang_choice):
    if lang_choice == '繁體中文':
        return easyocr.Reader(['ch_tra', 'en'])
    else:
        return easyocr.Reader(['ch_sim', 'en'])

reader = load_ocr(lang_option)

uploaded_file = st.file_uploader("請上傳您的 PPTX 檔案", type="pptx")

if uploaded_file:
    if st.button("🚀 開始強力去文字"):
        with st.spinner('正在掃描圖片並清除所有殘留符號...'):
            prs = Presentation(uploaded_file)
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == 13: # 圖片類型
                        try:
                            img_bytes = shape.image.blob
                            nparr = np.frombuffer(img_bytes, np.uint8)
                            img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                            if img is None: continue
                            
                            results = reader.readtext(img)
                            
                            # 建立初始遮罩
                            mask = np.zeros(img.shape[:2], dtype=np.uint8)
                            found_any = False
                            
                            for (bbox, text, prob) in results:
                                found_any = True
                                (tl, tr, br, bl) = bbox
                                pts = np.array([tl, tr, br, bl], np.int32)
                                cv2.fillPoly(mask, [pts], 255)
                            
                            if found_any:
                                # 【關鍵修正：遮罩膨脹】
                                # 建立一個圓形內核，將遮罩向外擴大，吃掉邊緣標點
                                kernel = np.ones((clean_strength, clean_strength), np.uint8)
                                mask = cv2.dilate(mask, kernel, iterations=1)
                                
                                # 執行影像修補
                                dst = cv2.inpaint(img, mask, 7, cv2.INPAINT_TELEA)
                                
                                _, buffer = cv2.imencode(".png", dst)
                                new_img_stream = io.BytesIO(buffer)
                                
                                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                                slide.shapes.add_picture(new_img_stream, left, top, width, height)
                                
                                # 移除舊圖
                                pic_element = shape._element
                                pic_element.getparent().remove(pic_element)
                                
                        except Exception as e:
                            continue

            output = io.BytesIO()
            prs.save(output)
            st.success("✅ 強化清除完成！")
            st.download_button(
                label="📥 下載強化處理後的 PPTX",
                data=output.getvalue(),
                file_name="cleaned_pro_file.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
